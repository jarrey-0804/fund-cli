"""
PPT报告生成器.

使用 python-pptx 生成 PowerPoint 格式的基金分析报告。
"""

import tempfile
from pathlib import Path
from typing import Any

from fund_cli.core.reporter import Reporter


class PptxReporter(Reporter):
    """PPT报告生成器."""

    def generate(
        self,
        fund_code: str,
        metrics: dict[str, Any],
        nav_data: Any = None,
        benchmark_data: Any = None,
        **kwargs,
    ) -> str:  # type: ignore[override]
        """生成PPT内容（返回临时文件路径）."""
        try:
            from pptx import Presentation
            from pptx.dml.color import RGBColor
            from pptx.enum.text import PP_ALIGN
            from pptx.util import Inches, Pt
        except ImportError as exc:
            raise RuntimeError("python-pptx 未安装，请运行: pip install python-pptx") from exc

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        # 幻灯片1：封面
        slide_layout = prs.slide_layouts[6]  # blank
        slide = prs.slides.add_slide(slide_layout)

        # 标题
        txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11.333), Inches(2))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = f"{fund_code} 基金分析报告"
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0x1A, 0x52, 0x76)
        p.alignment = PP_ALIGN.CENTER

        # 日期
        from datetime import date

        p2 = tf.add_paragraph()
        p2.text = date.today().strftime("%Y年%m月%d日")
        p2.font.size = Pt(20)
        p2.font.color.rgb = RGBColor(0x7F, 0x8C, 0x8D)
        p2.alignment = PP_ALIGN.CENTER

        # 幻灯片2：核心指标
        slide2 = prs.slides.add_slide(slide_layout)

        txBox2 = slide2.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
        tf2 = txBox2.text_frame
        p3 = tf2.paragraphs[0]
        p3.text = "核心绩效指标"
        p3.font.size = Pt(32)
        p3.font.bold = True
        p3.font.color.rgb = RGBColor(0x2C, 0x3E, 0x50)

        # 指标表格
        rows, cols = 6, 4
        table_shape = slide2.shapes.add_table(
            rows, cols, Inches(0.5), Inches(1.5), Inches(12), Inches(4)
        )
        table = table_shape.table

        headers = ["指标", "值", "指标", "值"]
        for i, h in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = h
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(14)
                paragraph.font.bold = True

        metrics_pairs = [
            (
                "总收益率",
                metrics.get("total_return", "N/A"),
                "年化收益率",
                metrics.get("annualized_return", "N/A"),
            ),
            (
                "波动率",
                metrics.get("volatility", "N/A"),
                "夏普比率",
                metrics.get("sharpe_ratio", "N/A"),
            ),
            ("最大回撤", metrics.get("max_drawdown", "N/A"), "Alpha", metrics.get("alpha", "N/A")),
        ]
        for r, (n1, v1, n2, v2) in enumerate(metrics_pairs, 1):
            for c, val in enumerate([n1, v1, n2, v2]):
                cell = table.cell(r, c)
                if isinstance(val, (int, float)):
                    cell.text = f"{val:.4f}"
                else:
                    cell.text = str(val)
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.size = Pt(13)

        # 幻灯片3：投资建议
        slide3 = prs.slides.add_slide(slide_layout)
        txBox3 = slide3.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
        tf3 = txBox3.text_frame
        p4 = tf3.paragraphs[0]
        p4.text = "投资摘要"
        p4.font.size = Pt(32)
        p4.font.bold = True

        summary_box = slide3.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(4))
        stf = summary_box.text_frame
        stf.word_wrap = True
        sp = stf.paragraphs[0]
        total_return = metrics.get("total_return", 0)
        sp.text = f"总收益率 {total_return:.2%}"
        sp.font.size = Pt(18)
        sp.font.color.rgb = (
            RGBColor(0x27, 0xAE, 0x60) if total_return > 0 else RGBColor(0xE7, 0x4C, 0x3C)
        )

        sp2 = stf.add_paragraph()
        sp2.text = f"夏普比率 {metrics.get('sharpe_ratio', 'N/A')}，最大回撤 {metrics.get('max_drawdown', 'N/A')}"
        sp2.font.size = Pt(16)
        sp2.space_before = Pt(12)

        # 保存到临时文件
        temp_path = Path(tempfile.gettempdir()) / f"{fund_code}_report.pptx"
        prs.save(str(temp_path))
        return str(temp_path)

    def save(self, content: str, output_path: str) -> None:
        """保存PPT文件."""
        import shutil

        shutil.copy2(content, output_path)

    def get_formats(self) -> list[str]:
        return ["pptx"]
